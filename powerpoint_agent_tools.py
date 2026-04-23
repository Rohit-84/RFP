import os
import io
import time
import requests
import re
from typing import List, Dict, Any, Optional
# import logging
from logger_config import get_logger
import psutil
import traceback
import math
from urllib.parse import quote
from concurrent.futures import ThreadPoolExecutor
from threading import Lock  #  ADDED: Import Lock
from dotenv import load_dotenv
from thefuzz import fuzz  # Requires: pip install thefuzz
from azure.core.credentials import AzureKeyCredential
from azure.search.documents import SearchClient
from openai import AzureOpenAI
import pythoncom
import win32com.client
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE

load_dotenv()

# --- Logging Setup ---
# logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
# logger = logging.getLogger(__name__)
logger = get_logger(__name__)

FOLDER_ID = os.getenv("FOLDER_ID")
SHAREPOINTURL = os.getenv("SHAREPOINTURL")
G_CAPABILITIES_PREFIX = "01 - PRE-SALES HANDY/G. Capabilities Deck"

#Added: Lock (traffic control) for powerpoint generation
ppt_generation_lock = Lock()

# ------------------------------------
# 1. CONFIGURATION & BLACKLISTS
# ------------------------------------

GENERIC_BLACKLIST_KEYWORDS = [
    "l&t group", "world’s largest", "world's largest", "defense",
    "nuclear", "training", "pod", "odc", "staffing", "engagement model",
    "resource", "talent", "construction", "heavy engineering", "infrastructure",
    "factory", "plant", "boiler", "turbine", "reactor", "metro", "stadium"
]

SEARCH_FLUFF_KEYWORDS = [
    "revenue", "financial", "headquarters", "locations", "awards", "glance",
    "founded", "employees", "offices", "nasdaq", "stock", "presence", "gartner"
]

CAPABILITY_KEYWORDS = {
    "Tableau": ["tableau", "viz", "workbook", "hyper", "dashboard", "story", "calculated field", "visualization"],
    "AI ML": ["ai", "genai", "machine learning", "mlops", "llm", "rag", "agent", "deep learning", "nlp", "computer vision", "transformer", "neural network", "generative", "inference", "hugging face", "openai", "llama"],
    "RPA": ["rpa", "uipath", "automation", "robot", "bot", "orchestrator", "automation anywhere", "blue prism", "power automate", "unattended"],
    "Cloud Services": ["azure", "gcp", "aws", "cloud", "migration", "modernization", "devops", "kubernetes", "docker", "serverless", "terraform", "microservices", "landing zone"],
    "Databricks": ["databricks", "spark", "lakehouse", "delta", "mlflow", "unity catalog", "medallion", "notebook"],
    "Data Management & Analytics": ["cdp", "data architecture", "data engineering", "etl", "elt", "pipeline", "warehouse", "governance", "analytics", "business intelligence"],
    "Healthcare": ["healthcare", "hipaa", "fhir", "clinical", "ehr", "hl7", "interoperability", "patient", "provider", "payer"],
    "Salesforce": ["salesforce", "crm", "service cloud", "sales cloud", "apex", "lightning", "mulesoft", "force.com"],
    "Snowflake": ["snowflake", "data warehouse", "snowpipe", "data sharing", "etl", "elt", "snowpark", "cortex"],
    "DEFAULT": []
}

CAPABILITY_EXCLUSIONS = {
    "Tableau": ["hypermesh", "ansys", "electronics"],
    "AI ML": ["hypermesh", "ansys", "manual testing", "sap", "infra", "snowflake", "salesforce", "tableau"],
    "RPA": ["generative ai", "llm", "data warehouse"],
    "Cloud Services": ["on-premise", "mainframe", "pixel", "hardware", "carrier"],
    "Databricks": ["snowflake", "synapse", "redshift", "power bi"],
    "Healthcare": ["finance", "retail", "manufacturing", "automotive"],
    "Salesforce": ["dynamics", "hubspot", "zoho", "sap crm"],
    "Snowflake": ["databricks", "delta lake", "hadoop"],
    "DEFAULT": []
}

#new text extraction function
#fetches from all possible shapes/elements
def _slide_text(slide):
    """
    Extract and normalize all readable text from a PowerPoint slide.

    This function walks through:
      1) Text frames (text boxes, placeholders, etc.)
      2) Tables (iterating through all cells)
      3) Grouped shapes (recursively)
    and collects their text content.

    The extracted text fragments are normalized using `_normalize_text` to ensure:
      - Lowercasing
      - Collapsing whitespace
      - Removing the standalone word 'agenda'
      - Trimming

    Parameters
    ----------
    slide : pptx.slide.Slide
        The PowerPoint slide object to extract text from.

    Returns
    -------
    str
        A single normalized string containing all slide text.
    """
    parts = []

    def extract_from_shape(shape):
        # 1) Text frames (regular text boxes, placeholders, etc.)
        if getattr(shape, "has_text_frame", False) and getattr(shape, "text_frame", None):
            txt = shape.text_frame.text or ""
            norm = _normalize_text(txt)
            if norm:
                parts.append(norm)

        # 2) Tables (shape_type == 19 corresponds to MSO_SHAPE_TYPE.TABLE)
        if getattr(shape, "shape_type", None) == 19:
            try:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        cell_txt = cell.text or ""
                        norm = _normalize_text(cell_txt)
                        if norm:
                            parts.append(norm)
            except Exception:
                pass

        # 3) SmartArt/groups treated as grouped shapes (shape_type == 6 is MSO_SHAPE_TYPE.GROUP)
        if getattr(shape, "shape_type", None) == 6:
            try:
                for shp in shape.shapes:
                    extract_from_shape(shp)
            except Exception:
                pass

    # Iterate through all shapes on the slide
    for shape in slide.shapes:
        try:
            extract_from_shape(shape)
        except Exception:
            continue

    # Join and normalize the final output once more for safety
    return _normalize_text(" ".join(parts))


def _normalize_text(t: str | None) -> str:
    """
    Normalize a text string for consistent downstream processing.

    This function:
      1) Converts text to lowercase.
      2) Collapses all runs of whitespace (spaces, tabs, newlines) into single spaces.
      3) Removes the standalone word 'agenda' (case-insensitive).
      4) Trims leading and trailing whitespace.

    Parameters
    ----------
    t : str | None
        The input text. If None is provided, it will be treated as an empty string.

    Returns
    -------
    str
        A clean, simplified, and normalized string.

    Examples
    --------
    >>> _normalize_text("   Agenda:\\n  Project  PLAN   ")
    'project plan'

    >>> _normalize_text(None)
    ''
    """
    t = (t or "").lower()
    t = re.sub(r"\s+", " ", t)
    t = re.sub(r"\b(agenda)\b", "", t)
    return t.strip()


def _signature_for_slide(slide) -> str:
    """
    Build a normalized text 'signature' for a slide.

    The signature is used as a compact, comparable representation of a slide’s content.
    Steps:
      1) Extract all text from the slide via `_slide_text(slide)`.
      2) Normalize it using `_normalize_text`.
      3) Truncate to the first 1000 characters for lightweight comparisons or indexing.

    Parameters
    ----------
    slide : Any
        A slide-like object compatible with `_slide_text(slide)`.

    Returns
    -------
    str
        A normalized text snippet (up to 1000 characters) representing the slide.

    Notes
    -----
    - This is useful for deduplication, similarity checks, or quick indexing.
    """
    txt = _normalize_text(_slide_text(slide))
    return txt[:1000]


def _looks_like_thank_you(slide) -> bool:
    """
    Heuristically detect if a slide appears to be a 'Thank You' slide.

    Logic:
      1) Extract full slide text via `_slide_text(slide)`.
      2) Lowercase and collapse whitespace.
      3) Search for the phrase 'thank you' allowing flexible spacing (e.g., 'thank    you').

    Parameters
    ----------
    slide : Any
        A slide-like object compatible with `_slide_text(slide)`.

    Returns
    -------
    bool
        True if the slide text contains 'thank you' as a phrase; otherwise False.

    Examples
    --------
    - Matches: 'THANK    you', 'thank you!', 'We want to say THANK YOU to everyone'
    - Non-matches: 'thanks', 'with thanks', 'thankyour team' (no word boundary)
    """
    txt = (_slide_text(slide) or "").lower()
    txt = re.sub(r"\s+", " ", txt)
    return bool(re.search(r"\bthank\s*you\b", txt))

def _is_title_like_slide(slide):
    # Counters for different types of text-bearing placeholders/shapes
    title_placeholders = 0
    body_placeholders = 0
    other_text_shapes = 0

    try:
        # Extract normalized full-slide text (already cleaned by _slide_text → _normalize_text)
        total_text = _slide_text(slide)

        # Iterate through all shapes on the slide
        for shape in slide.shapes:
            has_text = False

            # Check if shape has a text frame and contains non-empty text
            if getattr(shape, "has_text_frame", False) and shape.text_frame:
                txt = shape.text_frame.text or ""
                has_text = bool(txt.strip())

            # Check if the shape is a placeholder
            is_placeholder = hasattr(shape, 'is_placeholder') and shape.is_placeholder

            if is_placeholder:
                try:
                    # Identify the placeholder type
                    ph_type = shape.placeholder_format.type

                    # Title-type placeholders
                    if ph_type in (
                        PP_PLACEHOLDER_TYPE.TITLE,
                        PP_PLACEHOLDER_TYPE.CENTER_TITLE,
                        PP_PLACEHOLDER_TYPE.SUBTITLE,
                        PP_PLACEHOLDER_TYPE.VERTICAL_TITLE,
                    ):
                        if has_text:
                            title_placeholders += 1

                    # Body/content-type placeholders
                    elif ph_type in (
                        PP_PLACEHOLDER_TYPE.BODY,
                        PP_PLACEHOLDER_TYPE.CONTENT,
                        PP_PLACEHOLDER_TYPE.OBJECT,
                    ):
                        if has_text:
                            body_placeholders += 1

                    # Other placeholders with text
                    elif has_text:
                        other_text_shapes += 1

                except Exception:
                    # Fallback: if placeholder type lookup fails but text exists, treat as other text
                    if has_text:
                        other_text_shapes += 1

            # Non-placeholder shapes with text
            elif has_text:
                other_text_shapes += 1

        # Remove leading/trailing whitespace after extraction
        # Remove leading/trailing whitespace after extraction
        total_text = total_text.strip()

        # --- HEURISTICS FOR TITLE-LIKE SLIDES ---

        # 1) Slide has title placeholder(s) but no body/other text → typical title slide
        if title_placeholders >= 1 and body_placeholders == 0 and other_text_shapes == 0:
            return True

        # 2) Short overall text and no body placeholders → resembles title-only slide
        if total_text and len(total_text) <= 70 and body_placeholders == 0:
            return True

        # 3) Only one text-bearing element and overall text is short → likely title slide
        if (title_placeholders + body_placeholders + other_text_shapes == 1) and len(total_text) < 100:
            return True

    except Exception as e:
        # Log any issues encountered during detection
        logger.warning(f"Error checking slide title-likeness: {e}")

    # Default: not a title-like slide
    return False

def _kill_powerpoint():
    # Iterate through all running processes, requesting pid and name info
    for proc in psutil.process_iter(['pid', 'name']):
       
        # Check if process name exists and matches PowerPoint (POWERPNT.EXE)
        if proc.info['name'] and 'POWERPNT' in proc.info['name'].upper():
            try:
                # Attempt to terminate the process
                proc.kill()
            except:
                # Ignore errors such as access denied or process already closed
                pass

def _delete_slide_by_index(prs, index):
    try:
        # Ensure index is within valid slide range
        if 0 <= index < len(prs.slides):

            # Remove the slide's entry from the slide ID list
            # (python-pptx internally tracks slides in _sldIdLst)
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[index])

    except:
        # Ignore errors (e.g., invalid index, corruption, access issues)
        pass

def _move_slide(prs, old_idx, new_idx):
    try:
        # Validate that both old and new indices are within valid slide list bounds
        if 0 <= old_idx < len(prs.slides) and 0 <= new_idx <= len(prs.slides):

            # Internal list of slide IDs maintained by python-pptx
            sldIdLst = prs.slides._sldIdLst

            # Extract the slide element at the original index
            el = sldIdLst[old_idx]

            # Remove the slide from its old position
            sldIdLst.remove(el)

            # Insert the slide into its new position
            sldIdLst.insert(new_idx, el)

    except:
        # Silently ignore any errors (invalid index, internal errors, etc.)
        pass

# ------------------------------------
# 3. ROBUST CLEANUP
# ------------------------------------

def _cleanup_presentation(final_path):
    # Log the start of the cleanup with the file's base name
    logger.info(f"Starting cleanup process for: {os.path.basename(final_path)}")
    try:
        # Abort if the target file path doesn't exist
        if not os.path.exists(final_path): return False

        # Load the presentation from disk
        prs = Presentation(final_path)
        initial_slide_count = len(prs.slides)

        # If there are no slides, nothing to clean—consider it successful
        if initial_slide_count == 0: return True

        # the below code is shifted to find_relevant_slides()
        # ---------------------------------------------------------------------
        # Pass 1 (COMMENTED OUT): Remove title-like slides
        # NOTE:
        #   This block is intentionally disabled due to its tendency to remove
        #   important slides (especially in Tableau exports).
        #   If you want this functionality, improve `_is_title_like_slide()`
        #   and then re-enable this block.
        #
        # old function
        # # Pass 1: Remove title-like
        # title_slides_to_remove = []
        # for i, s in enumerate(prs.slides):
        #     is_last_slide = (i == initial_slide_count - 1)
        #     try:
        #         if _is_title_like_slide(s) and not _looks_like_thank_you(s) and not is_last_slide:
        #             title_slides_to_remove.append(i)
        #     except Exception: pass
        #
        # for i in sorted(title_slides_to_remove, reverse=True):
        #     _delete_slide_by_index(prs, i)
        # ---------------------------------------------------------------------

        # Pass 2: Deduplicate slides based on a signature
        # current_slides = prs.slides
        # seen_signatures = []   # store signatures of unique slides encountered
        # duplicate_indices = [] # indices of slides identified as duplicates

        # for i, s in enumerate(current_slides):
            # sig = _signature_for_slide(s)  # compute a textual/structural signature

            # # If signature is empty or too short, treat as duplicate (likely low-info slide)
            # if not sig or len(sig) < 20:
            #      duplicate_indices.append(i); continue

            # # Compare against seen signatures using fuzz.ratio; 99+ is considered a dup
            # is_dup = any(fuzz.ratio(sig, prev) >= 99 for prev in seen_signatures)
            # if is_dup:
            #     duplicate_indices.append(i)
            # else:
            #     seen_signatures.append(sig)
       
        # Remove duplicates starting from the end to keep indices stable
        # for i in sorted(duplicate_indices, reverse=True):
        #     _delete_slide_by_index(prs, i)

        # # Pass 3: Remove "Thank You" slides (common end slides)
        # current_slides = prs.slides
        # ty_positions = [i for i, s in enumerate(current_slides) if _looks_like_thank_you(s)]

        # # Delete detected "Thank You" slides, iterating in reverse index order
        # for i in sorted(ty_positions, reverse=True):
        #     _delete_slide_by_index(prs, i)

        # Save the cleaned presentation back to the same path
        prs.save(final_path)

        # Log success with the resulting slide count
        logger.info(f"✅ Cleanup finished. Final slide count: {len(prs.slides)}")
        return True

    except Exception as e:
        # Log the error and indicate cleanup failure
        logger.error(f"Cleanup failed: {e}")
        return False

# ------------------------------------
# 4. DOWNLOADER & CACHING
# ------------------------------------

def _download_worker(file_info, cache_dir, token, drive_id):
    """
    Downloads a file using:
    1) FAST direct path approach
    2) ONLY if that fails → SLOW search fallback
    """
    file_name = file_info.get("file_name")
    folder_path = file_info.get("folder_path")

    if not file_name:
        logger.error(f"Incomplete file_info: {file_info}")
        return

    local_path = os.path.join(cache_dir, file_name)

    if os.path.exists(local_path):
        logger.debug(f"File already cached: {local_path}")
        return

    headers = {"Authorization": f"Bearer {token}"}

    def _download(download_url):
        os.makedirs(os.path.dirname(local_path), exist_ok=True)

        with requests.get(download_url, stream=True, timeout=180) as r:
            r.raise_for_status()
            with open(local_path, "wb") as f:
                for chunk in r.iter_content(chunk_size=8192):
                    if chunk:
                        f.write(chunk)

    # ============================================================
    # ⚡ FAST APPROACH (STRICT)
    # ============================================================
    try:
        if not folder_path:
            raise ValueError("Missing folder_path for fast approach")

        full_graph_path = f"{folder_path}/{file_name}"
        encoded_path = quote(full_graph_path)

        url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/root:/{encoded_path}"

        meta_resp = requests.get(
            url,
            headers=headers,
            params={"$select": "@microsoft.graph.downloadUrl"},
            timeout=30
        )
        meta_resp.raise_for_status()

        download_url = meta_resp.json().get("@microsoft.graph.downloadUrl")

        if not download_url:
            raise ValueError("downloadUrl missing in fast approach")

        _download(download_url)

        logger.info(f"✅ FAST download succeeded: {file_name}")
        return  # IMPORTANT: stop here if fast succeeds

    except Exception as e:
        logger.warning(f"FAST approach failed for '{file_name}' Error '{e}'")

    # ============================================================
    # 🐢 SLOW APPROACH (ONLY if FAST failed)
    # ============================================================
    try:
        search_url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/root/search(q='{file_name}')"

        search_resp = requests.get(search_url, headers=headers, timeout=30)
        search_resp.raise_for_status()

        items = search_resp.json().get("value", [])

        file_item = next(
            (it for it in items if it.get("name") == file_name and "file" in it),
            None
        )

        if not file_item:
            raise ValueError(f"File '{file_name}' not found in search")

        download_url = file_item.get("@microsoft.graph.downloadUrl")

        if not download_url:
            item_id = file_item.get("id")
            if not item_id:
                raise ValueError("Missing item id in search result")

            meta_url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/items/{item_id}"
            meta_resp = requests.get(meta_url, headers=headers, timeout=30)
            meta_resp.raise_for_status()

            download_url = meta_resp.json().get("@microsoft.graph.downloadUrl")

        if not download_url:
            raise ValueError("downloadUrl still missing in slow approach")

        _download(download_url)

        logger.info(f"🐢 SLOW fallback succeeded: {file_name}")

    except Exception as e:
        logger.error(
            f"❌ BOTH approaches failed for '{file_name}'."
            f"FAST error | SLOW error: '{e}'",
            exc_info=True
        )

        # Cleanup partial file
        if os.path.exists(local_path):
            try:
                os.remove(local_path)
            except OSError:
                logger.warning(f"Could not delete partial file: {local_path}")

def _ensure_files_are_cached(files_to_check, cache_dir, token, drive_id):
    # Ensure the cache directory exists for storing downloaded files
    os.makedirs(cache_dir, exist_ok=True)

    # Determine which files are missing from the local cache
    # (Checks by file_name inside the cache directory)
    # missing = [
        # f for f in files_to_check
    # if not os.path.exists(os.path.join(cache_dir, files_to_check["file_name"])):
    #     return
    # ]

    # # If no files are missing, nothing to download
    # if not missing:
    #     return

    # Log how many files will be downloaded
    logger.info(f"Downloading {len(files_to_check)} files (Parallel)...")
    # logger.info(f"Downloading ")
    # _download_worker(file_info=files_to_check,cache_dir=cache_dir,drive_id=drive_id,token=token)
    # Use a thread pool to download missing files concurrently
    with ThreadPoolExecutor(max_workers=6) as executor:
        for f in files_to_check:
            # Submit each download task to the worker pool
            executor.submit(_download_worker(f, cache_dir, token, drive_id))

def _is_embedded_file(doc):
    # Extract the file name from the document metadata; default to empty string if missing
    fn = (doc.get("file_name") or "").strip()

    # Check if the file name ends with `.pptx` (case‑insensitive)
    # If not, this document is not considered an embedded PowerPoint file
    if not fn.lower().endswith(".pptx"):
        return False

    # If it is a `.pptx` file, treat it as an embedded file
    return True

# ------------------------------------
# 5. SMART SEARCH LOGIC (HYBRID + LLM)
# ------------------------------------

def _cosine_similarity(a, b):
    # Compute the dot product of the two vectors
    dot = sum(x * y for x, y in zip(a, b))

    # Compute the magnitude (L2 norm) of vector a
    na = math.sqrt(sum(x * x for x in a))

    # Compute the magnitude (L2 norm) of vector b
    nb = math.sqrt(sum(y * y for y in b))

    # If either vector has zero magnitude, similarity is undefined → return 0
    if na == 0 or nb == 0:
        return 0.0

    # Return cosine similarity = dot / (|a| * |b|)
    return dot / (na * nb)

def _expand_query_with_llm(user_query, capability, openai_client, chat_deployment, context_profile=None):
    """
    Uses chat model to expand the user's query into a more specific, technical search query.
    Falls back to original query on any failure.
    """
    # If no deployment/model is provided, return the original query unchanged
    if not chat_deployment:
        return user_query

    try:
        # Optional context segment appended to the prompt
        ctx = ""
        if context_profile:
            ctx = f"\nContext profile: {context_profile}\n"
            
        system_prompt = f"""
You are a search optimizer for a slide library.
Capability: {capability}
User query: "{user_query}"
{ctx}

Extract ONLY the contextual texts to search from the user query and return it concisely.
Return ONLY the extracted context text.
"""

        # Call the chat completion endpoint with low temperature for determinism
        resp = openai_client.chat.completions.create(
            model=chat_deployment,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": "Return the improved search query."}
            ],
            temperature=0,
            max_tokens=80,
        )

        # Extract the improved query text from the first choice
        improved = resp.choices[0].message.content.strip()

        # If model returned non-empty output, use it and log the expansion
        if improved:
            logger.info(f"🧠 LLM expanded query:\n  Original: {user_query}\n  Expanded: {improved}")
            return improved

        # Fallback to original query if response is empty
        return user_query

    except Exception as e:
        # On any failure (API, parsing, etc.), log a warning and return original query
        logger.warning(f"LLM query expansion failed, using original query. Error: {e}")
        return user_query

import threading

def find_relevant_slides(
    endpoint: str,
    key: str,
    index_name: str,
    openai_endpoint: str,
    openai_key: str,
    openai_deployment: str,
    chat_deployment: str,
    query: str,
    num_slides: int,
    cache_dir: str,
    token: str,
    drive_id: str,
    capability: Optional[str] = None,
    context_profile: Optional[str] = None,
    filter_mode: Optional[str] = None,
    target_files: Optional[List[str]] = None,
    exclude_files_content: Optional[List[str]] = None
) -> Dict[str, List[Any]]:
    """
    Returns a dict with two lists:
    - "exclude_files_content": list of excluded slide contents
    - "final_slides": list of dicts with file_name and slide_number
    """

    G_FOLDER_PREFIX = "01 - PRE-SALES HANDY/G. Capabilities Deck"

    try:
        search_client = SearchClient(endpoint, index_name, AzureKeyCredential(key))
        openai_client = AzureOpenAI(
            api_key=openai_key,
            api_version="2024-02-15-preview",
            azure_endpoint=openai_endpoint,
        )

        expanded_query = _expand_query_with_llm(
            user_query=query,
            capability=capability,
            openai_client=openai_client,
            chat_deployment=chat_deployment,
            context_profile=context_profile,
        )

        query_emb = openai_client.embeddings.create(
            input=[expanded_query],
            model=openai_deployment
        ).data[0].embedding

        num_slides_k =int(0)
        if num_slides == 0:
            num_slides_k = int(1)
        else:
            num_slides_k = num_slides

        results = search_client.search(
            search_text=expanded_query,
            vector_queries=[{
                "vector": query_emb,
                "k": int(num_slides_k * 4),
                "fields": "content_vector",
                "kind": "vector"
            }],
            select=["file_name", "folder_path", "slide_number", "slide_title", "content"],
            top=int(num_slides * 4),
            scoring_profile="boostcontext"
        )

        seen_content_in_file: List[str] = []
        seen_title_in_file: List[str] = []
        files_to_fetch: List[Dict[str, str]] = []
        candidates: List[Dict[str, Any]] = []
        exclude_files_content = exclude_files_content or []

        for r in results:
            fn = r.get("file_name")
            fp = r.get("folder_path")
            sn = r.get("slide_number")
            st = r.get("slide_title")
            cont = r.get("content")
            contx = str(r.get("slide_context"))

            if not fn or not fp:
                continue

            src = os.path.join(os.path.abspath(cache_dir), fn)
            if not os.path.exists(src):
                files_to_fetch.append({"folder_path": fp, "file_name": fn})

            if "This is a thank you slide." in contx:
                continue
            if "This is a title-only slide." in contx:
                continue
            if cont in seen_content_in_file or st in seen_title_in_file:
                continue
            if len(st) >= len(cont):
                continue
            if cont in exclude_files_content:
                continue

            exclude_files_content.append(cont)
            seen_content_in_file.append(cont)
            seen_title_in_file.append(st)

            candidates.append({"file_name": fn, "slide_number": sn})
            if len(candidates) == num_slides:
                seen_content_in_file.clear()
                seen_title_in_file.clear()
                exclude_files_contents = exclude_files_content.copy()
                exclude_files_content.clear()
                _ensure_files_are_cached(files_to_check=files_to_fetch, cache_dir=cache_dir, token=token, drive_id=drive_id)
                return {
                    "exclude_files_content": exclude_files_contents,
                    "final_slides": candidates
                }

        if not candidates:
            logger.warning("No candidate files returned from search.")
        seen_content_in_file.clear()
        seen_title_in_file.clear()
        exclude_files_contents = exclude_files_content.copy()
        exclude_files_content.clear()
        _ensure_files_are_cached(files_to_check=files_to_fetch, cache_dir=cache_dir, token=token, drive_id=drive_id)
        return {
            "exclude_files_content": exclude_files_content,
            "final_slides": candidates
        }

    except Exception as e:
        logger.error(f"AGENT TOOL ERROR in find_relevant_slides: {e}", exc_info=True)
        return {
            "exclude_files_content": [],
            "final_slides": []
        }
 

# ------------------------------------
# 6. DECK ASSEMBLY (ROBUST WIN32COM)
# ------------------------------------

def create_dynamic_pitch_deck(initial_slides, additional_slides, summary_text, output_path, cache_dir):
    # Start log for deck assembly
    logger.info("🛠️ Starting Deck Assembly...")
   
    # Resolve asset paths relative to this script's directory
    script_dir = os.path.dirname(os.path.abspath(__file__))
    BAR_IMG = os.path.join(script_dir, "Screenshot 2025-10-25 115739.png")  # Right-side brand bar
    LOGO_IMG = os.path.join(script_dir, "ltts_logo.png")                    # Bottom-left logo
    TY_SLIDE = os.path.join(script_dir, "ThankYouSlide.pptx")               # Optional "Thank You" deck
   
    # Absolute paths for output and cache
    output_path_abs = os.path.abspath(output_path)
    cache_dir_abs = os.path.abspath(cache_dir)

    with ppt_generation_lock:
        # Ensure no hanging PowerPoint processes interfere with COM automation
        _kill_powerpoint()
       
        # Placeholders for COM objects and state
        powerpoint = None
        master_pres = None
        opened_presentations = {}  # Cache of opened source presentations
        win32_success = False      # Tracks whether assembly succeeded

        try:
            # Initialize COM and spawn PowerPoint (visible, alerts suppressed)
            pythoncom.CoInitialize()
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = True
            powerpoint.DisplayAlerts = False
           
            # Create a new master presentation (destination)
            master_pres = powerpoint.Presentations.Add()
            opened_presentations["__master__"] = master_pres
           
            # Cache slide dimensions for positioning branding assets
            slide_width = master_pres.PageSetup.SlideWidth
            slide_height = master_pres.PageSetup.SlideHeight

            # Helper: Open a source presentation from cache (once) and reuse handle
            def get_source_pres(filename):
                if filename in opened_presentations: return opened_presentations[filename]
                try:
                    src = os.path.join(cache_dir_abs, filename)
                    if not os.path.exists(src): return None
                    p = powerpoint.Presentations.Open(src, ReadOnly=False, WithWindow=False)
                    opened_presentations[filename] = p
                    return p
                except: return None

            # Helper: Add branding assets (bar + logo) to a slide, if assets exist
            def add_branding(slide):
                try:
                    if os.path.exists(BAR_IMG):
                        bar_w = 14  # Fixed width for the right-side brand bar
                        pic = slide.Shapes.AddPicture(BAR_IMG, 0, 1, slide_width - bar_w, 0, bar_w, slide_height)
                        pic.ZOrder(1)  # Send backward to avoid overlaying slide content
                    if os.path.exists(LOGO_IMG):
                        w, h, m = 60, 25, 10  # Logo width, height, margin
                        slide.Shapes.AddPicture(LOGO_IMG, 0, 1, m, slide_height - h - m, w, h)
                except: pass

            # Helper: Copy a slide from a source file into the master deck safely
            def copy_slide_safe(filename, slide_idx):
                """
                Copy a single slide from a source presentation into the master deck.
                - Attempts clipboard copy/paste first (preferred method).
                - If clipboard fails, falls back to InsertFromFile.
                - Performs a corruption pre-check to skip unreadable files.
                """

                # Build absolute path to source presentation
                src_path = os.path.join(cache_dir_abs, filename)
                if not os.path.exists(src_path):
                    logger.warning(f"Missing source file: {src_path}")
                    return

                # --- Corruption pre-check ---
                try:
                    # Try opening and immediately closing the file to validate integrity
                    test_pres = powerpoint.Presentations.Open(src_path, ReadOnly=True, WithWindow=False)
                    total_slides = test_pres.Slides.Count  # Count slides for later validation
                    test_pres.Close()
                except Exception as e:
                    # If file cannot be opened, treat it as corrupted/unreadable
                    logger.error(f"File {filename} appears corrupted or unreadable: {e}")
                    return  # Skip this file entirely

                # --- Clipboard copy attempt ---
                try:
                    # Retrieve source presentation object
                    src_pres = get_source_pres(filename)
                    if not src_pres:
                        raise Exception("Source presentation not available")

                    # Validate slide index
                    if slide_idx < 1 or slide_idx > src_pres.Slides.Count:
                        logger.warning(f"Invalid slide index {slide_idx} in {filename} (has {src_pres.Slides.Count} slides)")
                        return

                    # Track number of slides before copy
                    cnt_before = master_pres.Slides.Count

                    # Copy slide to clipboard and paste into master deck
                    src_pres.Slides(slide_idx).Copy()
                    time.sleep(0.05)  # Small delay to allow clipboard to settle
                    master_pres.Slides.Paste(Index=cnt_before + 1)

                    # If paste succeeded, add branding to the newly added slide
                    if master_pres.Slides.Count > cnt_before:
                        add_branding(master_pres.Slides(master_pres.Slides.Count))

                    logger.info(f"Copied slide {slide_idx} from {filename} via clipboard")
                    return

                except Exception as e:
                    # Clipboard copy failed, log warning and proceed to fallback
                    logger.warning(f"Clipboard copy failed for {filename} slide {slide_idx}: {e}")

                # --- Fallback: InsertFromFile ---
                try:
                    # Validate slide index against total slides
                    if slide_idx < 1 or slide_idx > total_slides:
                        logger.warning(f"Invalid slide index {slide_idx} in {filename} (has {total_slides} slides)")
                        return

                    # Track number of slides before insertion
                    cnt_before = master_pres.Slides.Count

                    # Insert slide directly from file into master deck
                    master_pres.Slides.InsertFromFile(src_path, cnt_before, slide_idx, slide_idx)

                    # If insertion succeeded, add branding to the newly added slide
                    if master_pres.Slides.Count > cnt_before:
                        add_branding(master_pres.Slides(master_pres.Slides.Count))

                    logger.info(f"Inserted slide {slide_idx} from {filename} via InsertFromFile")

                except Exception as e:
                    # Both clipboard and InsertFromFile failed
                    logger.error(f"Failed to insert slide {slide_idx} from {filename}: {e}")
                   
            # Copy all requested slides: initial + additional
            for s in initial_slides + additional_slides:
                copy_slide_safe(s['file_name'], s['slide_number'])
               
            # Optionally append a standardized "Thank You" slide at the end
            if os.path.exists(TY_SLIDE):
                try:
                    ty_pres = powerpoint.Presentations.Open(
                        os.path.abspath(TY_SLIDE),
                        ReadOnly=False,
                        WithWindow=False
                    )
                    opened_presentations["__ty__"] = ty_pres

                    if ty_pres.Slides.Count >= 1:
                        try:
                            # Attempt to copy/paste the first slide
                            ty_pres.Slides(1).Copy()
                            master_pres.Slides.Paste()
                            logger.info("TY slide appended via Copy/Paste")
                        except Exception as copy_err:
                            logger.warning(f"Copy/Paste failed, falling back: {copy_err}")
                            try:
                                # Fallback: insert directly from file
                                master_pres.Slides.InsertFromFile(
                                    os.path.abspath(TY_SLIDE),
                                    master_pres.Slides.Count
                                )
                                logger.info("TY slide appended via InsertFromFile")
                            except Exception as insert_err:
                                logger.error(f"InsertFromFile also failed: {insert_err}")

                        # Add branding to the last slide if present
                        if master_pres.Slides.Count > 0:
                            add_branding(master_pres.Slides(master_pres.Slides.Count))

                except Exception as e:
                    logger.error(f"Failed to append TY slide: {e}")

            # Save the assembled presentation to the target output path
            master_pres.SaveAs(output_path_abs)
            win32_success = True
            logger.info("✅ Deck assembled and saved.")
           
        except Exception as e:
            # Log and bail on any fatal assembly errors
            logger.error(f"Assembly Fatal Error: {e}")
            return None
        finally:
            # Close all opened source presentations (except the master placeholder)
            for name, p in opened_presentations.items():
                if name == "__master__": continue
                try: p.Close()
                except: pass

            # Close the master presentation handle
            if master_pres:
                try: master_pres.Close()
                except: pass

            # Quit PowerPoint application if it was created
            if powerpoint:
                try: powerpoint.Quit()
                except: pass

            # Uninitialize COM context
            pythoncom.CoUninitialize()

    # Post-process the saved deck (e.g., dedupe slides, remove TY slides) on success
    if win32_success:
        _cleanup_presentation(output_path_abs)
       
    # Return the output path (or None if earlier failure)
    return output_path